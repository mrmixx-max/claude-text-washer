"""Document text extraction + metadata strip.

Adapted from TWS ``documents/service.py``, expanded to handle real document
formats: DOCX, XLSX, PPTX, EPUB, ODT, HTML, Markdown (and TXT).

Two responsibilities:
1. ``extract_text(data, filename) -> ExtractedDocument`` — pulls readable text
   out using the best available library (each imported lazily so the module
   loads without optional deps installed).
2. ``strip_metadata(data, filename) -> (bytes, report)`` — removes generator /
   AI-provenance metadata from the document (for formats that carry it).

Honest boundaries:
- OCR is NOT performed (no image text extraction inside documents).
- Lossy formats (XLSX cell layout, PPTX slide order) are reconstructed as
  flat text, not round-tripped back into the original part structure.
"""
from __future__ import annotations

import io
import re
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from pathlib import Path

# ---------------------------------------------------------------- data types
SUPPORTED = ["docx", "pdf", "xlsx", "pptx", "epub", "odt", "html", "htm", "md", "markdown", "txt"]


@dataclass
class ExtractedDocument:
    filename: str
    format: str
    text: str
    metadata: dict = field(default_factory=dict)

    def to_dict(self) -> dict:
        return {
            "filename": self.filename,
            "format": self.format,
            "chars": len(self.text),
            "metadata": self.metadata,
        }


@dataclass
class DocumentCleanReport:
    format: str
    actions: list[str] = field(default_factory=list)
    removed_keys: list[str] = field(default_factory=list)
    cleaned: bytes | None = None

    def to_dict(self) -> dict:
        return {
            "format": self.format,
            "actions": self.actions,
            "removed_keys": self.removed_keys,
        }


# ---------------------------------------------------------------- format detection
# magic-byte sniffing for container formats (all ZIP-based)
_DOCX_MAGIC = b"PK\x03\x04"


def detect_format(data: bytes, filename: str) -> str:
    """Detect a format, preferring magic bytes over extension."""
    if len(data) >= 4 and data[:4] == _DOCX_MAGIC:
        # Inspect the content to distinguish docx/xlsx/pptx/epub/odt
        try:
            zf = zipfile.ZipFile(io.BytesIO(data))
            names = set(zf.namelist())
        except zipfile.BadZipFile:
            ext = _ext(filename)
            return ext or "unknown"
        if "word/document.xml" in names:
            return "docx"
        if "xl/workbook.xml" in names:
            return "xlsx"
        if "ppt/slides/slide1.xml" in names or "ppt/presentation.xml" in names:
            return "pptx"
        if "mimetype" in names:
            try:
                if zf.read("mimetype").startswith(b"application/epub"):
                    return "epub"
            except (KeyError, RuntimeError):
                pass
            return "odt"
        ext = _ext(filename)
        return ext or "zip"
    if len(data) >= 5 and data[:5] == b"%PDF-":
        return "pdf"
    if data[:6] in (b"GIF87a", b"GIF89a"):
        return "gif"
    if len(data) >= 8 and data[:8] == b"\x89PNG\r\n\x1a\n":
        return "png"
    if len(data) >= 2 and data[:2] == b"\xff\xd8":
        return "jpeg"
    if len(data) >= 4 and data[:4] == b"RIFF" and len(data) >= 12 and data[8:12] == b"WEBP":
        return "webp"
    if len(data) >= 12 and data[4:8] == b"ftyp":
        brand = data[8:12]
        if brand in (b"avif", b"heic"):
            return brand.decode("ascii", "replace")
    if data[:2] in (b"II", b"MM"):
        return "tiff"
    text = data.lstrip()[:200].decode("utf-8", "ignore").lower()
    if text.startswith("<svg"):
        return "svg"
    if text.startswith("<html") or text.startswith("<!doctype html") or text.startswith("<head"):
        return "html"
    # fall back to extension
    ext = _ext(filename)
    return ext or "txt"


def _ext(filename: str) -> str:
    if not filename or "." not in filename:
        return ""
    return filename.rsplit(".", 1)[-1].lower()


# ---------------------------------------------------------------- text extraction
def extract_text(data: bytes, filename: str, *, force_format: str | None = None) -> ExtractedDocument:
    """Extract readable text from a document (any supported format).

    Parameters
    ----------
    data:
        Raw file bytes.
    filename:
        Filename (used for extension-based fallback detection).
    force_format:
        If given, bypass magic-byte detection and extract using this format
        key (e.g. ``"pdf"`` or ``"docx"``).  Unknown formats fall through to a
        best-effort UTF-8 decode.
    """
    fmt = force_format or detect_format(data, filename)
    metadata = {}
    text = ""

    if fmt == "docx":
        text, metadata = _extract_docx(data)
    elif fmt == "pdf":
        text, metadata = _extract_pdf(data)
    elif fmt == "xlsx":
        text, metadata = _extract_xlsx(data)
    elif fmt == "pptx":
        text, metadata = _extract_pptx(data)
    elif fmt == "epub":
        text, metadata = _extract_epub(data)
    elif fmt == "odt":
        text, metadata = _extract_odt(data)
    elif fmt in ("html", "htm"):
        text, metadata = _extract_html(data)
    elif fmt in ("md", "markdown"):
        text, metadata = _extract_markdown(data)
    elif fmt == "txt":
        text, metadata = _extract_txt(data)
    else:
        # best-effort: decode as utf-8 and return
        text = data.decode("utf-8", "replace")
        metadata = {"chars": len(text)}

    return ExtractedDocument(filename=filename, format=fmt, text=text, metadata=metadata)


def _extract_pdf(data: bytes) -> tuple[str, dict]:
    """Extract text from a PDF using PyMuPDF (``pymupdf``), lazily imported.

    Falls back to the deprecated ``fitz`` alias if the new import name is absent.
    Returns ``(text, metadata)`` where *metadata* holds the PDF's docinfo dict.
    """
    meta: dict = {}
    try:
        try:
            import pymupdf as fitz  # current recommended import name
        except ImportError:  # pragma: no cover - fallback path
            import fitz  # type: ignore  # legacy alias
    except ImportError:
        return data.decode("utf-8", "replace"), {"error": "pymupdf not installed"}

    texts: list[str] = []
    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:
        return data.decode("utf-8", "replace"), {"error": f"pymupdf open failed: {exc}"}
    try:
        for page in doc:
            texts.append(page.get_text())
        try:
            info = doc.metadata
            meta = {k: v for k, v in (info or {}).items() if v}
        except Exception:
            pass
    finally:
        doc.close()
    return "\n".join(texts).strip(), meta


def _extract_docx(data: bytes) -> tuple[str, dict]:
    # Lazy import — python-docx only needed for DOCX.
    from docx import Document  # type: ignore

    doc = Document(io.BytesIO(data))
    texts = [p.text for p in doc.paragraphs if p.text]
    meta = {}
    try:
        core = doc.core_properties
        meta = {
            "title": core.title,
            "author": core.author,
            "subject": core.subject,
            "created": str(core.created) if core.created else None,
            "modified": str(core.modified) if core.modified else None,
            "last_modified_by": core.last_modified_by,
        }
    except Exception:
        pass
    return "\n".join(texts), meta


def _extract_xlsx(data: bytes) -> tuple[str, dict]:
    # Lazy import — openpyxl only needed for XLSX.
    from openpyxl import load_workbook  # type: ignore

    wb = load_workbook(filename=io.BytesIO(data), read_only=True, data_only=True)
    texts: list[str] = []
    for ws in wb.worksheets:
        texts.append(ws.title)
        for row in ws.iter_rows(values_only=True):
            for val in row:
                if val is not None:
                    texts.append(str(val))
    wb.close()
    meta = {}
    try:
        cp = wb.properties
        meta = {"title": cp.title, "author": cp.creator, "subject": cp.subject}
    except Exception:
        pass
    return "\n".join(texts), meta


def _extract_pptx(data: bytes) -> tuple[str, dict]:
    # Lazy import — python-pptx only needed for PPTX.
    from pptx import Presentation  # type: ignore

    prs = Presentation(io.BytesIO(data))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts), {"slides": len(prs.slides)}


def _xml_text_iter(data: bytes, tag_suffix: str) -> list[str]:
    """Yield text from XML elements whose local tag ends with tag_suffix."""
    out: list[str] = []
    it = ET.iterparse(io.BytesIO(data))
    for _event, elem in it:
        if elem.tag.endswith(tag_suffix):
            if elem.text and elem.text.strip():
                out.append(elem.text.strip())
            for child in elem:
                if child.tail and child.tail.strip():
                    out.append(child.tail.strip())
        elem.clear()
    return out


def _extract_epub(data: bytes) -> tuple[str, dict]:
    """Extract text from an EPUB (ZIP of XHTML). Uses stdlib only."""
    texts = []
    meta = {}
    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile:
        return "", {}
    for name in zf.namelist():
        if name.endswith((".xhtml", ".html", ".htm", ".txt")):
            raw = zf.read(name)
            txt = _xml_text_iter(raw, "body") if name.endswith((".xhtml", ".html", ".htm")) else raw.decode("utf-8", "replace")
            if isinstance(txt, bytes):
                txt = txt.decode("utf-8", "replace")
            texts.append(txt)
    # container metadata from META-INF/container.xml / content.opf
    try:
        container = zf.read("META-INF/container.xml")
        for _e, _t in ET.iterparse(io.BytesIO(container)):
            pass
    except (KeyError, RuntimeError):
        pass
    return "\n\n".join(texts), meta


def _extract_odt(data: bytes) -> tuple[str, dict]:
    """Extract text from an ODT (OpenDocument) via stdlib XML."""
    texts = []
    meta = {}
    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile:
        return "", {}
    # text:content and text:s (spaces are represented as <text:s/>)
    for part_name in ("content.xml", "meta.xml"):
        if part_name in zf.namelist():
            raw = zf.read(part_name)
            try:
                root = ET.fromstring(raw)
            except ET.ParseError:
                continue
            ns = {"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0"}
            for para in root.iter("{urn:oasis:names:tc:opendocument:xmlns:text:1.0}p"):
                if para.text and para.text.strip():
                    texts.append(para.text.strip())
            for span in root.iter("{urn:oasis:names:tc:opendocument:xmlns:text:1.0}span"):
                if span.text:
                    texts.append(span.text)
    return "\n".join(texts), meta


def _extract_html(data: bytes) -> tuple[str, dict]:
    """Extract text from HTML. Prefers bs4, falls back to stdlib regex."""
    meta: dict = {}
    try:
        from bs4 import BeautifulSoup  # type: ignore

        soup = BeautifulSoup(data, "html.parser")
        meta = {name: content.get("content", "") for name, content in
                ((m.get("name"), m) for m in soup.find_all("meta")) if name}
        for script in soup(["script", "style", "noscript"]):
            script.decompose()
        return soup.get_text(separator="\n", strip=True), meta
    except ImportError:
        text = data.decode("utf-8", "ignore")
        text = re.sub(r"<script[\s\S]*?</script>", "", text, flags=re.IGNORECASE)
        text = re.sub(r"<style[\s\S]*?</style>", "", text, flags=re.IGNORECASE)
        text = re.sub(r"<noscript[\s\S]*?</noscript>", "", text, flags=re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
        # crude meta extraction
        for m in re.finditer(r'<meta\s+name=["\']([^"\']+)["\']\s+content=["\']([^"\']*)["\']>', text, re.IGNORECASE):
            meta[m.group(1)] = m.group(2)
        return text, meta


def _extract_markdown(data: bytes) -> tuple[str, dict]:
    meta: dict = {}
    text = data.decode("utf-8", "replace")
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            fm = text[3:end]
            for line in fm.splitlines():
                m = re.match(r"\s*([\w-]+)\s*:\s*(.*)", line)
                if m:
                    meta[m.group(1)] = m.group(2).strip()
            text = text[end + 5 :]
    return text.strip(), meta


def _extract_txt(data: bytes) -> tuple[str, dict]:
    from .encoding_detect import detect_encoding

    result = detect_encoding(data)
    source_enc = result.encoding.replace("-sig", "")
    text = data.decode(source_enc, errors="replace")
    if result.bom:
        # strip BOM bytes already handled by decode of utf-8-sig, but ensure
        _, bom_len = _bom_len(result.encoding, data)
        if bom_len:
            text = data[bom_len:].decode(source_enc, errors="replace")
    return text, {"encoding": result.encoding, "confidence": result.confidence}


def _bom_len(encoding: str, data: bytes) -> tuple[bool, int]:
    from .encoding_detect import detect_bom

    _enc, n = detect_bom(data)
    return _enc is not None, n


# ---------------------------------------------------------------- metadata strip
# ---------------------------------------------------------------- metadata strip for documents
def strip_metadata(data: bytes, filename: str) -> tuple[bytes, dict]:
    """Remove AI-provenance / generator metadata from a document."""
    fmt = detect_format(data, filename)
    if fmt == "docx":
        return _strip_docx_meta(data, filename)
    if fmt == "xlsx":
        return _strip_zip_meta(data, filename, xl_meta=True)
    if fmt == "pptx":
        return _strip_zip_meta(data, filename, pptx=True)
    if fmt in ("odt", "epub"):
        return _strip_zip_meta(data, filename, odt=fmt == "odt")
    if fmt in ("html", "htm"):
        return _strip_html_meta(data, filename)
    if fmt in ("md", "markdown"):
        return _strip_markdown_meta(data, filename)
    if fmt == "txt":
        cleaned, info = _strip_bom_bytes(data)
        rep = DocumentCleanReport(format="txt", actions=info.get("actions", []))
        rep.cleaned = cleaned
        return cleaned, rep.to_dict()
    # Unknown / image formats handled by image_metadata; nothing to do here.
    rep = DocumentCleanReport(format=fmt or "unknown", actions=["unsupported_format"])
    rep.cleaned = data
    return data, rep.to_dict()


def _strip_docx_meta(data: bytes, filename: str) -> tuple[bytes, dict]:
    """Strip core/app props + AI hints from a DOCX via its own cleaner + repair."""
    from .docx_repair import repair_docx, validate_docx

    rep = DocumentCleanReport(format="docx")
    try:
        zin = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile:
        rep.actions.append("not_a_docx_zip")
        rep.cleaned = data
        return data, rep.to_dict()

    names = zin.namelist()
    # Drop customXml + docProps, scrub core/app props of AI hints.
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            content = zin.read(item.filename)
            if item.filename.startswith("customXml/"):
                rep.removed_keys.append(item.filename)
                rep.actions.append("removed_customXml_part")
                continue
            if item.filename in ("docProps/core.xml", "docProps/app.xml"):
                content = _scrub_docx_props(content)
                rep.removed_keys.append(item.filename)
                rep.actions.append("scrubbed_docx_props")
            zout.writestr(item, content)
    cleaned = out.getvalue()

    # Repair structure so Word opens it cleanly after stripping.
    repair = repair_docx(cleaned)
    cleaned = repair.data or cleaned
    rep.actions.append("repaired_after_strip")
    rep.cleaned = cleaned
    return cleaned, rep.to_dict()


def _scrub_docx_props(content: bytes) -> bytes:
    text = content.decode("utf-8", "replace")
    for tag in (
        "dc:creator", "cp:lastModifiedBy", "cp:revision", "dc:title", "dc:subject",
        "dc:description", "dc:creator", "cp:manager", "cp:company", "cp:category",
        "cp:keywords", "w:docSupplier", "meta:generator",
    ):
        text = re.sub(rf"<{tag}[^>]*>[^<]*</{tag}>", f"<{tag}></{tag}>", text, flags=re.IGNORECASE)
    # Drop any XML element whose local name smells AI/generator/provenance.
    text = re.sub(
        r"<(?:generator|provenance|synthid|c2pa|content-credentials|"
        r"authoring-tool)[^>]*>[\s\S]*?</(?:generator|provenance|synthid|c2pa|"
        r"content-credentials|authoring-tool)>",
        "",
        text,
        flags=re.IGNORECASE,
    )
    return text.encode("utf-8")


def _strip_zip_meta(data: bytes, filename: str, xl_meta: bool = False, pptx: bool = False, odt: bool = False) -> tuple[bytes, dict]:
    """Generic ZIP-container metadata strip (XLSX/PPTX/ODT/EPUB)."""
    from .docx_repair import repair_docx

    fmt = detect_format(data, filename) or "docx"
    rep = DocumentCleanReport(format=fmt)
    try:
        zin = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile:
        rep.actions.append("not_a_zip")
        rep.cleaned = data
        return data, rep.to_dict()

    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            name = item.filename
            content = zin.read(item.filename)
            drop = False
            if name.startswith("customXml/"):
                drop = True
                rep.removed_keys.append(name)
                rep.actions.append("removed_customXml_part")
            elif odt and name == "meta.xml":
                content = _scrub_odt_meta(content)
                rep.removed_keys.append(name)
                rep.actions.append("scrubbed_odt_meta")
            elif name == "docProps/core.xml" or name == "docProps/app.xml":
                content = _scrub_docx_props(content)
                rep.removed_keys.append(name)
                rep.actions.append("scrubbed_doc_props")
            if not drop:
                zout.writestr(item, content)
    cleaned = out.getvalue()
    repair = repair_docx(cleaned)
    cleaned = repair.data or cleaned
    rep.actions.append("repaired_after_strip")
    rep.cleaned = cleaned
    return cleaned, rep.to_dict()


def _scrub_odt_meta(content: bytes) -> bytes:
    text = content.decode("utf-8", "replace")
    for tag in ("meta:initial-creator", "meta:creator", "meta:date", "meta:editing-cycles",
                "meta:editing-duration", "meta:generator", "dc:creator", "dc:date",
                "dc:title", "dc:subject", "dc:description", "dc:source", "meta:keyword"):
        text = re.sub(rf"<{tag}(\s[^>]*)?>[^<]*</{tag}>", f"<{tag}></{tag}>", text, flags=re.IGNORECASE)
    return text.encode("utf-8")


def _strip_html_meta(data: bytes, filename: str) -> tuple[bytes, dict]:
    """Remove AI-provenance meta tags and data-ai* attributes from HTML."""
    from .encoding_detect import detect_encoding

    rep = DocumentCleanReport(format="html")
    text = data.decode("utf-8", "replace")
    new = text

    def _meta_drop(m):
        tag = m.group(0)
        name_m = re.search(r'name=["\']([^"\']+)["\']', tag, re.IGNORECASE)
        if name_m and _AI_NAME_RE.search(name_m.group(1)):
            rep.removed_keys.append(name_m.group(1))
            rep.actions.append("removed_ai_meta_tag")
            return ""
        return tag

    new = re.sub(r"<meta\b[^>]*>", _meta_drop, new, flags=re.IGNORECASE)
    before = len(new)
    new = re.sub(r"<script\b[^>]*>application/ld\+json[\s\S]*?</script>", "", new, flags=re.IGNORECASE)
    if len(new) < before:
        rep.actions.append("removed_jsonld_provenance_block")
    new = re.sub(r"\sdata-ai[\w-]*=['\"][^\"']*['\"]", "", new, flags=re.IGNORECASE)
    if new != text:
        rep.removed_bytes = len(text.encode()) - len(new.encode())
    rep.cleaned = new.encode("utf-8")
    return rep.cleaned, rep.to_dict()


def _strip_markdown_meta(data: bytes, filename: str) -> tuple[bytes, dict]:
    """Remove AI-generator keys from YAML frontmatter of Markdown."""
    rep = DocumentCleanReport(format="markdown")
    text = data.decode("utf-8", "replace")
    new = text
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            fm = text[3:end]
            lines = fm.splitlines()
            kept = []
            for line in lines:
                m = re.match(r"\s*([\w-]+)\s*:", line)
                if m and _AI_KEY_HINTS_META.search(m.group(1)):
                    rep.removed_keys.append(m.group(1))
                    rep.actions.append("removed_ai_frontmatter_key")
                    continue
                kept.append(line)
            new = "---" + "\n".join(kept) + text[end:]
    if new != text:
        rep.removed_bytes = len(text.encode()) - len(new.encode())
    rep.cleaned = new.encode("utf-8")
    return rep.cleaned, rep.to_dict()


def _strip_bom_bytes(data: bytes) -> tuple[bytes, dict]:
    from .encoding_detect import detect_bom

    _enc, n = detect_bom(data)
    actions = []
    if n:
        actions.append("stripped_bom")
    return data[n:], {"actions": actions}


_AI_NAME_RE = re.compile(
    r"(generator|provenance|synthid|c2pa|content-credentials|"
    r"ai-(generated|assistant|model|tool|content)|created-by|"
    r"authoring-tool)",
    re.IGNORECASE,
)
_AI_KEY_HINTS_META = re.compile(
    r"(generated|generator|created|ai|model|prompt|provenance|c2pa|synthid)",
    re.IGNORECASE,
)
